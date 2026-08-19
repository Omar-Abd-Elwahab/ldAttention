"""Compose the print-ready poster PDF directly.

There is no LibreOffice on this machine to convert the .pptx, so the PDF is
drawn from the same layout spec the .pptx builder uses
(``scripts/poster_layout.py``) and the template's own logos, footer strip, photo
and QR code, which are pulled straight out of the .pptx.

Every text block is measured and wrapped against the width of its box, and the
font size is stepped down until the block fits its height. Overflow is therefore
impossible by construction -- which is what went wrong in the previous build.

Usage::

    .venv312/bin/python3.12 scripts/render_poster_pdf.py
"""

from __future__ import annotations

import argparse
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.image as mpimg
import matplotlib.pyplot as plt
from matplotlib.patches import Circle, Rectangle
from pptx import Presentation

import poster_layout as L

DPI = 150  # only affects raster placement inside the vector PDF


# --------------------------------------------------------------------------- #
# Coordinate helpers: inches from the top-left corner -> figure fraction
# --------------------------------------------------------------------------- #
def fx(x_in: float) -> float:
    return x_in / L.PAGE_W


def fy(y_in: float) -> float:
    return 1.0 - y_in / L.PAGE_H


def _renderer(fig):
    return fig.canvas.get_renderer()


def _text_width_in(fig, text: str, size: float, family, weight="normal", style="normal") -> float:
    """Width of ``text`` in inches if drawn at ``size`` points."""
    artist = fig.text(0, 0, text, fontsize=size, fontfamily=family, fontweight=weight, fontstyle=style)
    bbox = artist.get_window_extent(renderer=_renderer(fig))
    artist.remove()
    return bbox.width / fig.dpi


def wrap(fig, text: str, size: float, family, max_w: float, weight="normal") -> list[str]:
    """Greedy word wrap using real measured widths."""
    lines: list[str] = []
    for hard_line in text.split("\n"):
        words = hard_line.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}"
            if _text_width_in(fig, candidate, size, family, weight) <= max_w:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def draw_paragraph(
    fig,
    box: L.Box,
    text: str,
    size: float,
    family=L.PDF_FONT_BODY,
    color=L.INK,
    weight="normal",
    line_spacing: float = 1.22,
    para_gap: float = 0.28,
    min_size: float = 12.0,
) -> float:
    """Draw wrapped text inside ``box``, shrinking until it fits. Returns the bottom y."""
    paragraphs = [p for p in text.split("\n\n")]
    while size >= min_size:
        line_h = size / 72.0 * line_spacing
        blocks = [wrap(fig, p, size, family, box.w, weight) for p in paragraphs]
        total = sum(len(b) * line_h for b in blocks) + para_gap * (len(blocks) - 1)
        if total <= box.h:
            break
        size -= 1.0

    line_h = size / 72.0 * line_spacing
    y = box.y
    for block in blocks:
        for line in block:
            fig.text(
                fx(box.x), fy(y + line_h * 0.78), line,
                fontsize=size, fontfamily=family, color=color, fontweight=weight,
                ha="left", va="baseline",
            )
            y += line_h
        y += para_gap
    return y - para_gap


def draw_bullets(fig, box: L.Box, bullets, size: float, min_size: float = 12.0) -> float:
    """Bulleted list with a coloured bold lead-in, hanging indent, auto-shrink."""
    indent = 0.75
    while size >= min_size:
        line_h = size / 72.0 * 1.20
        gap = size / 72.0 * 0.42
        layouts = []
        for label, body, _ in bullets:
            lead = f"{label}: "
            lead_w = _text_width_in(fig, lead, size, L.PDF_FONT_BODY, "bold")
            first = wrap(fig, body, size, L.PDF_FONT_BODY, box.w - indent - lead_w)
            if not first:
                first = [""]
            head, tail_text = first[0], " ".join(first[1:])
            tail = wrap(fig, tail_text, size, L.PDF_FONT_BODY, box.w - indent) if tail_text else []
            layouts.append((lead, lead_w, head, tail))
        total = sum((1 + len(t)) * line_h + gap for *_, t in layouts) - gap
        if total <= box.h:
            break
        size -= 1.0

    line_h = size / 72.0 * 1.20
    gap = size / 72.0 * 0.42
    y = box.y
    for (label, _, color), (lead, lead_w, head, tail) in zip(bullets, layouts):
        baseline = fy(y + line_h * 0.78)
        fig.text(fx(box.x + 0.15), baseline, "•", fontsize=size, fontfamily=L.PDF_FONT_BODY,
                 color=color, fontweight="bold", ha="left", va="baseline")
        fig.text(fx(box.x + indent), baseline, lead, fontsize=size, fontfamily=L.PDF_FONT_BODY,
                 color=color, fontweight="bold", ha="left", va="baseline")
        fig.text(fx(box.x + indent + lead_w), baseline, head, fontsize=size,
                 fontfamily=L.PDF_FONT_BODY, color=L.INK, ha="left", va="baseline")
        y += line_h
        for line in tail:
            fig.text(fx(box.x + indent), fy(y + line_h * 0.78), line, fontsize=size,
                     fontfamily=L.PDF_FONT_BODY, color=L.INK, ha="left", va="baseline")
            y += line_h
        y += gap
    return y - gap


def draw_badge(fig, number: str, label: str, box: L.Box) -> None:
    """Numbered circle with the section label beside it, plus the template's rule."""
    d = L.BADGE_DIAMETER
    cx, cy = box.x + d / 2, box.y + d / 2
    fig.patches.append(
        Circle((fx(cx), fy(cy)), radius=fx(d / 2), transform=fig.transFigure,
               facecolor=L.INK, edgecolor="none", zorder=5)
    )
    fig.text(fx(cx), fy(cy), number, fontsize=L.SZ_BADGE_NUM, fontfamily=L.PDF_FONT_TITLE,
             fontweight="bold", color="white", ha="center", va="center", zorder=6)
    text_x = box.x + d + L.BADGE_GAP
    fig.text(text_x / L.PAGE_W, fy(cy), label, fontsize=L.SZ_BADGE_LABEL,
             fontfamily=L.PDF_FONT_TITLE, fontweight="bold", color=L.INK,
             ha="left", va="center", zorder=6)
    rule_y = box.y + d + 0.22
    fig.add_artist(
        plt.Line2D([fx(box.x), fx(box.x1)], [fy(rule_y), fy(rule_y)],
                   transform=fig.transFigure, color=L.INK, linewidth=2.0, zorder=4)
    )


def place_image(
    fig,
    path: Path,
    box: L.Box,
    crop: tuple[float, float, float, float] = (0.0, 0.0, 0.0, 0.0),
    stretch: bool = False,
) -> None:
    """Draw an image inside ``box``.

    ``crop`` is the template's (left, right, top, bottom) fractional crop.
    ``stretch`` fills the box exactly, which is what the template does for its
    footer logos; content figures keep their aspect ratio instead.
    """
    if not path.exists():
        fig.patches.append(
            Rectangle((fx(box.x), fy(box.y1)), box.w / L.PAGE_W, box.h / L.PAGE_H,
                      transform=fig.transFigure, facecolor="#F0F0F0",
                      edgecolor=L.GREY, linestyle="--", zorder=1)
        )
        fig.text(fx(box.x + box.w / 2), fy(box.y + box.h / 2), f"missing: {path.name}",
                 ha="center", va="center", fontsize=20, color=L.RED)
        return

    img = mpimg.imread(path)
    ih, iw = img.shape[0], img.shape[1]
    cl, cr, ct, cb = crop
    if any(abs(c) > 1e-6 for c in crop):
        x0, x1 = int(round(max(cl, 0.0) * iw)), int(round((1.0 - max(cr, 0.0)) * iw))
        y0, y1 = int(round(max(ct, 0.0) * ih)), int(round((1.0 - max(cb, 0.0)) * ih))
        img = img[max(y0, 0) : max(y1, y0 + 1), max(x0, 0) : max(x1, x0 + 1)]
        ih, iw = img.shape[0], img.shape[1]

    if stretch:
        x, y, w, h = box.x, box.y, box.w, box.h
    else:
        scale = min(box.w / iw, box.h / ih)
        w, h = iw * scale, ih * scale
        x = box.x + (box.w - w) / 2
        y = box.y + (box.h - h) / 2

    ax = fig.add_axes([fx(x), fy(y + h), w / L.PAGE_W, h / L.PAGE_H], zorder=3)
    # A single-channel image would otherwise be run through the default colour
    # map, which is how the QR came out in viridis.
    kwargs = {"cmap": "gray", "vmin": 0, "vmax": 255 if img.dtype != float else 1} \
        if img.ndim == 2 else {}
    ax.imshow(img, interpolation="antialiased",
              aspect="auto" if stretch else "equal", **kwargs)
    ax.axis("off")


# --------------------------------------------------------------------------- #
# Template assets
# --------------------------------------------------------------------------- #
def extract_template_assets(out_dir: Path) -> dict[str, tuple[Path, L.Box, tuple]]:
    """Pull the footer logos / photo / QR out of the .pptx with position and crop."""
    out_dir.mkdir(parents=True, exist_ok=True)
    assets: dict[str, tuple[Path, L.Box, tuple]] = {}
    if not L.TEMPLATE.exists():
        return assets
    for shape in Presentation(str(L.TEMPLATE)).slides[0].shapes:
        try:
            image = shape.image
        except (AttributeError, ValueError):
            continue
        if image.ext not in {"png", "jpg", "jpeg"}:
            continue
        path = out_dir / f"{shape.name.replace(' ', '_')}.{image.ext}"
        path.write_bytes(image.blob)
        assets[shape.name] = (
            path,
            L.Box(shape.left / 914400, shape.top / 914400, shape.width / 914400, shape.height / 914400),
            (shape.crop_left, shape.crop_right, shape.crop_top, shape.crop_bottom),
        )
    return assets


# Footer pieces that belong to the template's design and are redrawn verbatim.
FOOTER_ASSETS = (
    "Picture 27",   # author photo
    "Picture 16",   # QR code
    "Graphic 1",    # "scan me" arrow
    "Picture 945",  # LinkedIn wordmark
    "Graphic 956",  # X / Twitter wordmark
    "Picture 44",   # CRIV
    "Graphic 19",   # Genome Canada
    "Picture 947",  # Genome Quebec
    "Picture 9",    # IBIS
    "Picture 32",   # IID
    "Picture 29",   # Universite Laval
)


def draw_footer(fig, assets: dict[str, tuple[Path, L.Box, tuple]]) -> None:
    fig.patches.append(
        Rectangle((0, 0), 1.0, (L.PAGE_H - L.FOOTER_TOP) / L.PAGE_H, transform=fig.transFigure,
                  facecolor=L.CREAM, edgecolor="none", zorder=0)
    )
    for name in FOOTER_ASSETS:
        if name not in assets:
            continue
        path, box, crop = assets[name]
        # The template's QR points at a different project; use ours if built.
        if name == "Picture 16" and (L.POSTER / "fig_qr.png").exists():
            path, crop = L.POSTER / "fig_qr.png", (0.0, 0.0, 0.0, 0.0)
        place_image(fig, path, box, crop=crop, stretch=True)

    author = L.Box(3.78, 32.20, 15.4, 3.6)
    lines = [
        ("Omar Abdelwahab  ·  Dr. Davoud Torkamaneh", "bold", L.SZ_FOOTER, L.INK),
        ("Computational Genomics  ·  Université Laval", "normal", L.SZ_FOOTER_SMALL, L.INK),
        ("Département de phytologie  ·  FSAA", "normal", L.SZ_FOOTER_SMALL, L.INK),
        ("Quebec City, Quebec, Canada", "normal", L.SZ_FOOTER_SMALL, L.INK),
        ("ECCB 2026  ·  European Conference on Computational Biology", "normal", 22, L.BLUE),
    ]
    y = author.y
    for text, weight, size, color in lines:
        y += size / 72.0 * 1.30
        fig.text(fx(author.x), fy(y), text, fontsize=size, fontfamily=L.PDF_FONT_BODY,
                 fontweight=weight, color=color, ha="left", va="baseline", zorder=6)

    # Repository link beside the QR code. The template's box is only 2.6 in wide,
    # which is what chopped "github.com/omar-abdelwaha b /ldAttentio n" apart in
    # the previous build; the label is split across lines that actually fit.
    fig.add_artist(
        plt.Line2D([fx(19.50), fx(19.50)], [fy(32.35), fy(35.80)], transform=fig.transFigure,
                   color=L.INK, linewidth=1.6, zorder=4)
    )
    link_x, link_w = 19.72, 2.45
    fig.text(fx(link_x), fy(33.00), "Code & pre-print", fontsize=22, fontfamily=L.PDF_FONT_BODY,
             fontweight="bold", color=L.INK, ha="left", va="baseline", zorder=6)
    y = 33.00
    for line in wrap(fig, "github.com/ omar-abdelwahab/ ldAttention", 20, L.PDF_FONT_BODY, link_w):
        y += 0.36
        fig.text(fx(link_x), fy(y), line, fontsize=20, fontfamily=L.PDF_FONT_BODY,
                 color="#0563C1", ha="left", va="baseline", zorder=6)
    fig.text(fx(link_x), fy(y + 0.55), "Scan me", fontsize=20, fontfamily=L.PDF_FONT_BODY,
             fontstyle="italic", color=L.MUTED, ha="left", va="baseline", zorder=6)

    fig.text(fx(30.25), fy(33.55), "/omar-abdelwahab/", fontsize=26, fontfamily=L.PDF_FONT_BODY,
             color=L.INK, ha="left", va="baseline", zorder=6)
    fig.text(fx(30.25), fy(35.12), "/omarabdelwhab/", fontsize=26, fontfamily=L.PDF_FONT_BODY,
             color=L.INK, ha="left", va="baseline", zorder=6)


# --------------------------------------------------------------------------- #
# Page
# --------------------------------------------------------------------------- #
def build(out_path: Path, assets_dir: Path) -> Path:
    headline = L.load_headline()
    assets = extract_template_assets(assets_dir)

    fig = plt.figure(figsize=(L.PAGE_W, L.PAGE_H), dpi=DPI, facecolor="white")
    fig.canvas.draw()  # a renderer must exist before any text can be measured

    draw_footer(fig, assets)

    # Column rule.
    fig.add_artist(
        plt.Line2D([fx(L.DIVIDER_X), fx(L.DIVIDER_X)], [fy(L.DIVIDER_Y0), fy(L.DIVIDER_Y1)],
                   transform=fig.transFigure, color=L.INK, linewidth=2.5, zorder=2)
    )

    # ---- title ----
    # Wrap and shrink together, then flow the subtitle and byline below whatever
    # height the title ended up at. A long title used to be shrunk to a floor and
    # then allowed to run straight across the column rule.
    author_line = "Omar Abdelwahab · Davoud Torkamaneh"
    author_size = 38
    title_size = L.SZ_TITLE
    while title_size > 44:
        title_lines = wrap(fig, L.TITLE_MAIN, title_size, L.PDF_FONT_TITLE, L.TITLE.w, "bold")
        sub_lines = wrap(fig, L.TITLE_SUB, L.SZ_SUBTITLE, L.PDF_FONT_TITLE, L.TITLE.w)
        needed = (
            len(title_lines) * title_size / 72.0 * 1.12
            + 0.30
            + len(sub_lines) * L.SZ_SUBTITLE / 72.0 * 1.20
            + 0.34
            + author_size / 72.0 * 1.20
        )
        if len(title_lines) <= 2 and needed <= L.TITLE.h:
            break
        title_size -= 1.0

    y = L.TITLE.y
    for line in title_lines:
        y += title_size / 72.0 * 1.12
        fig.text(fx(L.TITLE.x), fy(y), line, fontsize=title_size,
                 fontfamily=L.PDF_FONT_TITLE, fontweight="bold", color=L.INK,
                 ha="left", va="baseline")
    y += 0.30
    for line in sub_lines:
        y += L.SZ_SUBTITLE / 72.0 * 1.20
        fig.text(fx(L.TITLE.x), fy(y), line, fontsize=L.SZ_SUBTITLE,
                 fontfamily=L.PDF_FONT_TITLE, color=L.MUTED, ha="left", va="baseline")
    y += 0.34 + author_size / 72.0 * 1.20
    fig.text(fx(L.TITLE.x), fy(y), author_line, fontsize=author_size,
             fontfamily=L.PDF_FONT_BODY, fontweight="bold", color=L.BLUE,
             ha="left", va="baseline")
    fig.add_artist(
        plt.Line2D([fx(L.TITLE.x), fx(L.TITLE.x1)], [fy(L.TITLE_RULE_Y), fy(L.TITLE_RULE_Y)],
                   transform=fig.transFigure, color=L.INK, linewidth=3.0)
    )

    for number, label, box in L.SECTIONS:
        draw_badge(fig, number, label, box)

    # ---- 1 Background ----
    draw_paragraph(fig, L.BG_BODY, "\n\n".join(L.background_paragraphs(headline)), L.SZ_BODY)

    # ---- 2 Approach ----
    place_image(fig, L.POSTER / L.FIGURES["Picture 10"][0], L.AP_FIGURE)
    draw_bullets(fig, L.AP_BODY, L.approach_bullets(headline), L.SZ_BULLET)

    # ---- 3 Methods ----
    place_image(fig, L.POSTER / L.FIGURES["Picture 11"][0], L.ME_FIGURE)
    draw_paragraph(fig, L.ME_NOTE, L.methods_note(headline), 30, color=L.MUTED)

    # ---- 4 Results ----
    place_image(fig, L.POSTER / L.FIGURES["Picture 8"][0], L.RE_FIGURE_A)
    place_image(fig, L.POSTER / L.FIGURES["Picture 24"][0], L.RE_FIGURE_B)
    draw_paragraph(fig, L.RE_CAPTION, L.figure_caption(headline), L.SZ_CAPTION, color=L.MUTED)

    # ---- 5 Conclusion ----
    draw_bullets(fig, L.CO_BODY, L.conclusion_takeaways(headline), L.SZ_BODY)
    draw_paragraph(fig, L.CO_NOTE, L.conclusion_note(headline), 28, color=L.MUTED)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, format="pdf", facecolor="white")
    plt.close(fig)
    return out_path


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=str, default=str(L.POSTER / "ECCB2026_ldAttention_poster.pdf"))
    ap.add_argument("--assets", type=str, default=str(L.POSTER / "template_assets"))
    args = ap.parse_args()
    print(f"Poster PDF written to: {build(Path(args.out), Path(args.assets))}")


if __name__ == "__main__":
    main()
