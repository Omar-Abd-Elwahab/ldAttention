"""Build the editable ECCB 2026 poster (PPTX) from the IID template.

Keeps the template's visual identity — palette, serif title, numbered section
badges, cream footer band with logos, photo and QR — and replaces the content.
Geometry and copy both come from ``scripts/poster_layout.py``, the same spec the
print-PDF renderer uses, so the two outputs cannot drift apart.

The template's own section markers are grouped shapes whose label sits in a box
narrower than the word it holds, which is what produced "Backgrou nd" and
"Metho" in the previous build. They are dropped and redrawn here as a circle
plus a label on one wide line.
"""

from __future__ import annotations

import argparse
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import poster_layout as L

OUT = L.POSTER / "ECCB2026_ldAttention_poster.pptx"

# Template section-marker groups, replaced by drawn badges.
MARKER_GROUPS = ("Group 118", "Group 4", "Group 95", "Group 18", "Group 128")


def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


def place(shape, box: L.Box) -> None:
    shape.left, shape.top = Inches(box.x), Inches(box.y)
    shape.width, shape.height = Inches(box.w), Inches(box.h)


def drop(shape) -> None:
    shape._element.getparent().remove(shape._element)


def write_paragraphs(shape, lines, align=PP_ALIGN.JUSTIFY, space_after=14) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (text, bold, size, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        run.font.name = L.PPTX_FONT_BODY
        run.font.bold = bold
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = rgb(color)


def write_title(shape) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    rows = (
        (L.TITLE_MAIN, L.SZ_TITLE, True, L.PPTX_FONT_TITLE, L.INK),
        (L.TITLE_SUB, L.SZ_SUBTITLE, False, L.PPTX_FONT_TITLE, L.MUTED),
        ("Omar Abdelwahab · Davoud Torkamaneh", 38, True, L.PPTX_FONT_BODY, L.BLUE),
    )
    for i, (line, size, bold, font, color) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)


def write_bullets(shape, bullets) -> None:
    from pptx.oxml import parse_xml

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (label, body, color) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(14)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "685800")
        pPr.set("indent", "-685800")
        pPr.insert(0, parse_xml(
            '<a:buChar char="•" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        ))
        head = p.add_run()
        head.text = f"{label}: "
        head.font.bold = True
        head.font.name = L.PPTX_FONT_BODY
        head.font.size = Pt(L.SZ_BULLET)
        head.font.color.rgb = rgb(color)
        tail = p.add_run()
        tail.text = body
        tail.font.name = L.PPTX_FONT_BODY
        tail.font.size = Pt(L.SZ_BULLET)


def add_badge(slide, number: str, label: str, box: L.Box) -> None:
    """Numbered circle with the section label beside it, on one line."""
    d = L.BADGE_DIAMETER
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(box.x), Inches(box.y), Inches(d), Inches(d)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(L.INK)
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = number
    run.font.name = L.PPTX_FONT_TITLE
    run.font.bold = True
    run.font.size = Pt(L.SZ_BADGE_NUM)
    run.font.color.rgb = rgb("FFFFFF")

    text_x = box.x + d + L.BADGE_GAP
    caption = slide.shapes.add_textbox(
        Inches(text_x), Inches(box.y - 0.10), Inches(box.x1 - text_x), Inches(d + 0.20)
    )
    ctf = caption.text_frame
    ctf.word_wrap = False
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.LEFT
    crun = cp.add_run()
    crun.text = label
    crun.font.name = L.PPTX_FONT_TITLE
    crun.font.bold = True
    crun.font.size = Pt(L.SZ_BADGE_LABEL)
    crun.font.color.rgb = rgb(L.INK)

    rule_y = box.y + d + 0.22
    line = slide.shapes.add_connector(
        1, Inches(box.x), Inches(rule_y), Inches(box.x1), Inches(rule_y)
    )
    line.line.color.rgb = rgb(L.INK)
    line.line.width = Pt(2.0)


def replace_picture(slide, shape, image_path: Path, box: L.Box):
    """Swap a template picture for a generated figure, fitted inside its box."""
    drop(shape)
    from PIL import Image

    with Image.open(image_path) as im:
        iw, ih = im.size
    scale = min(box.w / iw, box.h / ih)
    w, h = iw * scale, ih * scale
    return slide.shapes.add_picture(
        str(image_path),
        Inches(box.x + (box.w - w) / 2),
        Inches(box.y + (box.h - h) / 2),
        width=Inches(w),
        height=Inches(h),
    )


def build(out_path: Path = OUT) -> Path:
    L.POSTER.mkdir(parents=True, exist_ok=True)
    if not L.TEMPLATE.exists():
        raise FileNotFoundError(f"IID template not found: {L.TEMPLATE}")
    missing = [n for n, (fn, _) in L.FIGURES.items() if not (L.POSTER / fn).exists()]
    if missing:
        raise FileNotFoundError(
            f"Missing figures for slots {missing}; run scripts/make_poster_figures.py first"
        )

    headline = L.load_headline()
    shutil.copy2(L.TEMPLATE, out_path)
    prs = Presentation(str(out_path))
    slide = prs.slides[0]

    for shape in list(slide.shapes):
        name = shape.name or ""
        text = shape.text_frame.text.strip() if shape.has_text_frame else ""

        if name in MARKER_GROUPS:
            drop(shape)

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and name in L.FIGURES:
            filename, box = L.FIGURES[name]
            replace_picture(slide, shape, L.POSTER / filename, box)

        # The template's QR points at a different project.
        elif name == "Picture 16" and (L.POSTER / "fig_qr.png").exists():
            box = L.Box(shape.left / 914400, shape.top / 914400,
                        shape.width / 914400, shape.height / 914400)
            replace_picture(slide, shape, L.POSTER / "fig_qr.png", box)

        elif name == "Title 1" or text.startswith("Refinement"):
            write_title(shape)
            place(shape, L.TITLE)

        elif name == "Straight Connector 15":
            shape.left, shape.top = Inches(L.LEFT_X0), Inches(L.TITLE_RULE_Y)
            shape.width, shape.height = Inches(L.LEFT_W), Emu(0)

        elif name == "TextBox 12" or text.startswith("Genetic variants"):
            write_paragraphs(
                shape,
                [(p, False, L.SZ_BODY, None) for p in L.background_paragraphs(headline)],
            )
            place(shape, L.BG_BODY)

        # The template splits Background across two boxes; one is enough here.
        elif name == "TextBox 28" or text.startswith("However, raw"):
            drop(shape)

        elif name == "TextBox 13" or "Heuristic filtering" in text:
            write_bullets(shape, L.approach_bullets(headline))
            place(shape, L.AP_BODY)

        elif name == "TextBox 26" or text.startswith("Figure1") or text.startswith("Figure "):
            write_paragraphs(shape, [(L.figure_caption(headline), False, L.SZ_CAPTION, L.MUTED)])
            place(shape, L.RE_CAPTION)

        elif name == "TextBox 132" or text.startswith("This study introduces"):
            write_bullets(shape, L.conclusion_takeaways(headline))
            place(shape, L.CO_BODY)

        elif name == "TextBox 137" or text == "Pre-print":
            write_paragraphs(
                shape,
                [("Code & pre-print", True, 22, L.INK), (L.REPO_URL, False, 20, "0563C1")],
                align=PP_ALIGN.LEFT,
                space_after=4,
            )

        elif "Omar Abdelwahab" in text:
            write_paragraphs(
                shape,
                [
                    ("Omar Abdelwahab  ·  Dr. Davoud Torkamaneh", True, L.SZ_FOOTER, None),
                    ("Computational Genomics  ·  Université Laval", False, L.SZ_FOOTER_SMALL, None),
                    ("Département de phytologie  ·  FSAA", False, L.SZ_FOOTER_SMALL, None),
                    ("Quebec City, Quebec, Canada", False, L.SZ_FOOTER_SMALL, None),
                    ("ECCB 2026  ·  European Conference on Computational Biology", False, 22, L.BLUE),
                ],
                align=PP_ALIGN.LEFT,
                space_after=2,
            )

    # Methods caption sits under the protocol strip; the template has no such box.
    note = slide.shapes.add_textbox(
        Inches(L.ME_NOTE.x), Inches(L.ME_NOTE.y), Inches(L.ME_NOTE.w), Inches(L.ME_NOTE.h)
    )
    write_paragraphs(note, [(L.methods_note(headline), False, 28, L.MUTED)])

    # Scope caveat under the conclusion takeaways; also not in the template.
    caveat = slide.shapes.add_textbox(
        Inches(L.CO_NOTE.x), Inches(L.CO_NOTE.y), Inches(L.CO_NOTE.w), Inches(L.CO_NOTE.h)
    )
    write_paragraphs(caveat, [(L.conclusion_note(headline), False, 28, L.MUTED)])

    for number, label, box in L.SECTIONS:
        add_badge(slide, number, label, box)

    prs.core_properties.title = "LD-Aware Attention — ECCB 2026"
    prs.core_properties.subject = "ldAttention"
    prs.core_properties.author = "Omar Abdelwahab"
    prs.core_properties.keywords = "linkage disequilibrium, attention, transformer, genomics, imputation"
    prs.save(str(out_path))
    return out_path


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=str, default=str(OUT))
    args = ap.parse_args()
    print(f"Poster saved to: {build(Path(args.out))}")


if __name__ == "__main__":
    main()
